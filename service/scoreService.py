import nltk
from nltk import word_tokenize
from nltk.corpus import stopwords
from nltk.stem.wordnet import WordNetLemmatizer
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
# ppt parsing
import os
from pptx import Presentation
import string
import re
import numpy as np
import pandas as pd

stop_words = stopwords.words('English') + list(string.printable)

lemmatizer = WordNetLemmatizer()

def clean_texts(w): 
    """Words pre-processing include removing regex, stopwords and tokenization."""
    w = re.sub(r'([^\s\w]|_)+', ' ', str(w)).lower()
    #w = str(w).lower()
    tokenize = word_tokenize(w)
    cleaned_tokens = [lemmatizer.lemmatize(word) for word in tokenize if word not in stop_words]
    return ' '.join(cleaned_tokens)


def extract_texts(filename): 
    """Extract texts from slides"""
    ppt_path = os.path('path'+filename)
    ppt = Presentation(ppt_path)
    text_list = []
    for slide_num, slide in enumerate(ppt.slides):
        extract_text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                extract_text += shape.text + "\n"
        text_list.append(extract_text)
    return text_list

def getTextDataframe(text_list):
    cleanTexts = list(map(clean_texts,text_list))
    text_df = pd.DataFrame({'cleaned_texts': cleanTexts})
    return text_df

def calculate_similarity(ppt_texts, topic_descs):
    # 将 PPT 文本与 topic 简介合并成一个列表
    combined_texts = ppt_texts.tolist() + topic_descs.tolist()
    
    # 计算 TF-IDF 矩阵
    tfidf_model = TfidfVectorizer()
    tfidf_matrix = tfidf_model.fit_transform(combined_texts)
    
    # 分割 TF-IDF 矩阵为 PPT 文本和 topic 简介的子矩阵
    ppt_tfidf = tfidf_matrix[:len(ppt_texts)]
    topic_tfidf = tfidf_matrix[len(ppt_texts):]
    
    # 计算每个 PPT 文本与每个 topic 简介之间的余弦相似度
    cosine_sim_matrix = cosine_similarity(ppt_tfidf, topic_tfidf)
    
    # 将相似度矩阵转换为 DataFrame
    cosine_sim_df = pd.DataFrame(cosine_sim_matrix, index=ppt_texts.index, columns=topic_descs.index)
    
    return cosine_sim_df