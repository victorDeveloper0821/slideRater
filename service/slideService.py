from pptx import Presentation
from database import db
from models import Slide, BulletPoint, Submission

## need to be checked
def extract_pptx_content(pptx_path, submission_id):
    """extract texts in slides into data table"""
    prs = Presentation(pptx_path)
    
    for slide_number, slide in enumerate(prs.slides, start=1):
        # Create a new Slide entry
        slide_data = Slide(
            submission_id=submission_id,
            slide_number=slide_number,
            title="",
            footer="",
            content=""
        )

        for shape in slide.shapes:
            # Check if shape has text frame
            if not shape.has_text_frame:
                continue

            # Extract title
            if shape == slide.shapes.title:
                slide_data.title = shape.text

            # Extract footer (assuming footers are identified by name)
            if "Footer" in shape.name:
                slide_data.footer = shape.text

            # Extract content and bullet points
            for paragraph in shape.text_frame.paragraphs:
                # If it's a bullet point (i.e., has a level)
                if paragraph.level > 0:
                    bullet_point = BulletPoint(
                        level=paragraph.level,
                        text=paragraph.text
                    )
                    slide_data.bullet_points.append(bullet_point)
                else:
                    # Otherwise, it's general content
                    slide_data.content += paragraph.text + "\n"
        
        # Add the Slide (and its bullet points) to the database session
        db.session.add(slide_data)
    
    # Commit all the changes to the database
    db.session.commit()


if __name__ == '__main__':
# 使用示例
    pptx_path = 'your_pptx_file.pptx'
    ppt_data = extract_pptx_content(pptx_path)

# 顯示提取的數據
    for slide in ppt_data:
        print(f"Slide {slide['slide_number']} Title: {slide['title']}")
        print(f"Footer: {slide['footer']}")
        print("Content:")
        for text in slide['content']:
            print(f"  - {text}")
            print("Bullet Points:")
        for bullet in slide['bullet_points']:
            print(f"  Level {bullet['level']}: {bullet['text']}")
        print("\n")
